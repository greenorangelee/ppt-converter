import io
import json
import os
import anthropic
from pptx import Presentation


# ─── Template info ────────────────────────────────────────────────────────────

def get_template_info(path: str) -> dict:
    prs = Presentation(path)
    layouts = []
    for i, layout in enumerate(prs.slide_layouts):
        phs = [
            {"idx": ph.placeholder_format.idx, "name": ph.name}
            for ph in layout.placeholders
        ]
        layouts.append({"index": i, "name": layout.name, "placeholders": phs})

    return {
        "filename": os.path.basename(path),
        "slide_count": len(prs.slides),
        "layout_count": len(prs.slide_layouts),
        "layouts": layouts,
    }


# ─── Content extraction ───────────────────────────────────────────────────────

def _extract_slide_content(prs: Presentation) -> list[dict]:
    slides = []
    for i, slide in enumerate(prs.slides):
        data: dict = {"index": i, "title": "", "body": "", "extra_texts": []}

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text = shape.text_frame.text.strip()
            if not text:
                continue

            if shape.is_placeholder:
                idx = shape.placeholder_format.idx
                if idx == 0:
                    data["title"] = text
                elif idx == 1:
                    lines = []
                    for para in shape.text_frame.paragraphs:
                        t = para.text.strip()
                        if t:
                            prefix = "  " * para.level if para.level else ""
                            lines.append(f"{prefix}{t}")
                    data["body"] = "\n".join(lines)
                else:
                    data["extra_texts"].append(text)
            else:
                data["extra_texts"].append(text)

        # fallback: if no title placeholder, use first extra text
        if not data["title"] and data["extra_texts"]:
            data["title"] = data["extra_texts"].pop(0)

        slides.append(data)
    return slides


# ─── Claude mapping ───────────────────────────────────────────────────────────

async def _ask_claude_mapping(api_key: str, layouts: list, slides: list) -> list:
    layout_desc = "\n".join(
        f"  [{l['index']}] \"{l['name']}\"  placeholders: {[p['name'] for p in l['placeholders']]}"
        for l in layouts
    )
    slides_json = json.dumps(slides, ensure_ascii=False, indent=2)

    prompt = f"""You are converting a PowerPoint presentation to match a company template.

Available layouts in the template:
{layout_desc}

User slides (JSON):
{slides_json}

Task: For each user slide, choose the best fitting layout and map slide content to placeholder indices.

Return ONLY a valid JSON array — no markdown, no explanation:
[
  {{
    "slide_index": 0,
    "layout_index": <integer>,
    "placeholders": {{
      "0": "<title text>",
      "1": "<body text, use \\n for bullet separation>"
    }}
  }}
]

Rules:
- Use layout index whose name contains "Title Slide" or "Cover" for an opening/cover slide.
- Use a "Title and Content" style layout for regular content slides.
- Preserve ALL text from the original slide; merge extra_texts into body if needed.
- Do not invent content.
- placeholder key "0" = title, "1" = main body/content."""

    client = anthropic.AsyncAnthropic(api_key=api_key)
    message = await client.messages.create(
        model="claude-sonnet-4-6",
        max_tokens=4096,
        messages=[{"role": "user", "content": prompt}],
    )

    raw = message.content[0].text.strip()
    # Strip optional markdown code fences
    if "```" in raw:
        raw = raw.split("```")[1]
        if raw.startswith("json"):
            raw = raw[4:]
        raw = raw.split("```")[0]

    return json.loads(raw.strip())


# ─── Build output PPT ─────────────────────────────────────────────────────────

def _build_ppt(template_path: str, mapping: list, output_path: str):
    # Load template into memory buffer so original file stays clean
    with open(template_path, "rb") as f:
        buf = io.BytesIO(f.read())
    prs = Presentation(buf)

    # Remove all existing slides (preserves slide masters / layouts)
    sldIdLst = prs.slides._sldIdLst
    for sldId in list(sldIdLst):
        sldIdLst.remove(sldId)

    max_layout = len(prs.slide_layouts) - 1

    for slide_map in mapping:
        layout_idx = max(0, min(int(slide_map.get("layout_index", 1)), max_layout))
        layout = prs.slide_layouts[layout_idx]
        slide = prs.slides.add_slide(layout)

        ph_content: dict = slide_map.get("placeholders", {})

        for ph in slide.placeholders:
            key = str(ph.placeholder_format.idx)
            text = ph_content.get(key, "")
            if not text:
                continue

            tf = ph.text_frame
            tf.clear()
            lines = text.split("\n")
            for i, line in enumerate(lines):
                if not line.strip():
                    continue
                if i == 0:
                    tf.paragraphs[0].text = line.strip()
                else:
                    para = tf.add_paragraph()
                    para.text = line.strip()

    prs.save(output_path)


# ─── Public entry point ───────────────────────────────────────────────────────

async def convert_ppt(template_path: str, input_path: str, output_path: str, api_key: str):
    template_info = get_template_info(template_path)
    layouts = template_info["layouts"]

    user_prs = Presentation(input_path)
    slides = _extract_slide_content(user_prs)

    if not slides:
        raise ValueError("변환할 슬라이드가 없습니다.")

    mapping = await _ask_claude_mapping(api_key, layouts, slides)
    _build_ppt(template_path, mapping, output_path)
